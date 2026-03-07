import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

async def extract_text_from_file(content: bytes, filename: str) -> str:
    """
    Extracts text from various file formats based on filename extension.
    """
    extension = filename.lower().split('.')[-1]
    
    try:
        if extension == 'pdf':
            return extract_from_pdf(content)
        elif extension == 'docx':
            return extract_from_docx(content)
        elif extension == 'pptx':
            return extract_from_pptx(content)
        else:
            # Default to text/plain
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                # If it's not UTF-8, it might be a binary file we don't support
                logger.error(f"Failed to decode file {filename} as UTF-8")
                raise ValueError(f"Unsupported or corrupted file format: {filename}")
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}")
        if isinstance(e, ValueError):
            raise e
        raise ValueError(f"Could not extract text from {filename}: {str(e)}")

def extract_from_pdf(content: bytes) -> str:
    pdf_file = io.BytesIO(content)
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_from_docx(content: bytes) -> str:
    docx_file = io.BytesIO(content)
    doc = Document(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_from_pptx(content: bytes) -> str:
    pptx_file = io.BytesIO(content)
    prs = Presentation(pptx_file)
    text_runs = []
    for Slide in prs.slides:
        for shape in Slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)
